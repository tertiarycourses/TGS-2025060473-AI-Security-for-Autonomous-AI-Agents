#!/usr/bin/env python3
"""AI Security for Autonomous AI Agents — trainer slide deck (TGS-2025060473).

Component library adapted from the house COMPACT v2 reference
(.claude/skills/wsq-slides/reference/compact/build_slides.py).

House rules enforced here:
  · two trainer profile cards (general template + Dr Alfred Ang)
  · Download Course Material as a browser-mock visual, never a bare link
  · Assessment Flow diagram
  · Briefing BEFORE Assessment
  · TRAQOM digital attendance at the front AND the end
  · closing block: Assessment → Assessment Flow → Digital Attendance → Thank You
  · NO practice exam (non-certification course — its absence is correct)
  · NO step-by-step procedures (those live in the Learner Guide)
"""

import os, sys, json, math, re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
import deck_content as D

ASSETS = os.path.join(HERE, "assets")
OUTDIR = HERE

# ---------------- palette ----------------
BLUE = RGBColor(0x1F, 0x6F, 0xEB); TEAL = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B); INK = RGBColor(0x16, 0x1B, 0x26)
GREY = RGBColor(0x5B, 0x63, 0x72); LIGHT = RGBColor(0xF5, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); LINE = RGBColor(0xE2, 0xE8, 0xF0)
VIOLET = RGBColor(0x7C, 0x3A, 0xED); RED = RGBColor(0xDC, 0x26, 0x26)
NAVY = RGBColor(0x0B, 0x12, 0x20); CODEBLUE = RGBColor(0x9C, 0xDC, 0xFE)
PALETTE = [BLUE, TEAL, VIOLET, AMBER]
CMAP = {"BLUE": BLUE, "TEAL": TEAL, "VIOLET": VIOLET, "AMBER": AMBER, "RED": RED, "GREY": GREY}

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(): return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(1, x, y, w, h); sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp


def oval(s, x, y, w, h, color):
    sp = s.shapes.add_shape(9, x, y, w, h); sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False; return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for t, sz, col, bold in ln:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col; r.font.name = "Arial"
    return tb


def bullets(s, x, y, w, h, items, size=18, color=INK, gap=10, mcolor=BLUE):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(gap)
        lvl = it[1] if isinstance(it, tuple) else 0
        text = it[0] if isinstance(it, tuple) else it
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + text
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = color if lvl == 0 else GREY
        r.font.name = "Arial"
        r.font.bold = (lvl == 0 and isinstance(it, tuple) and len(it) > 2 and it[2])
    return tb


PAGE = {"n": 0}; SLIDE_MAP = {}; CURRENT_META = {"evidence": "", "sources": []}


def mark(key): SLIDE_MAP[key] = PAGE["n"] + 1


def _domain(url):
    """Short human-readable host for a source URL, e.g. 'klarna.com'."""
    m = re.sub(r"^https?://(www\.)?", "", url or "")
    return m.split("/")[0] if m else ""


def _source_band(s, cite):
    """A visible, clickable 'Sources' strip for real-life case-study slides so
    learners can see and check the origin of every claim on the slide. Sits in the
    footer zone (below any slide caption/note, above the org/copyright line) so it
    never collides with card or table notes."""
    y = Inches(6.80); x = Inches(0.85); w = Inches(11.63); h = Inches(0.22)
    rect(s, x, y, w, h, LIGHT); rect(s, x, y, Inches(0.08), h, RED)
    # Build one run per source: "Publisher (domain)", hyperlinked to the URL.
    tb = s.shapes.add_textbox(x + Inches(0.22), y, w - Inches(0.4), h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    lead = p.add_run(); lead.text = "Sources:  "
    lead.font.size = Pt(9); lead.font.bold = True; lead.font.name = "Arial"
    lead.font.color.rgb = RED
    for i, src in enumerate(cite):
        # Take the publisher part before an em-dash if present, else the domain.
        title = src.get("title", "")
        publisher = title.split(" — ")[0].strip() if " — " in title else title
        if len(publisher) > 46:
            publisher = publisher[:43].rstrip() + "..."
        dom = _domain(src.get("url", ""))
        label = f"{publisher} ({dom})" if dom else publisher
        if i:
            sep = p.add_run(); sep.text = "   ·   "
            sep.font.size = Pt(9); sep.font.name = "Arial"; sep.font.color.rgb = GREY
        r = p.add_run(); r.text = label
        r.font.size = Pt(9); r.font.name = "Arial"; r.font.color.rgb = BLUE
        if src.get("url"):
            try:
                r.hyperlink.address = src["url"]
            except Exception:
                pass


def footer(s):
    PAGE["n"] += 1
    evidence = CURRENT_META.get("evidence", "")
    source_ids = ", ".join(CURRENT_META.get("sources", []))
    cite = CURRENT_META.get("cite", [])
    if cite:
        # The visible Sources band replaces the cryptic "Sources: S.." ID line for
        # these slides (it occupies the same footer zone), so we do not draw both.
        _source_band(s, cite)
    elif evidence or source_ids:
        source_line = "Evidence: " + evidence if evidence else ""
        if source_ids:
            source_line += ("  ·  " if source_line else "") + "Sources: " + source_ids
        txt(s, Inches(0.85), Inches(6.82), Inches(11.5), Inches(0.18),
            [[(source_line, 6.8, GREY, False)]], align=PP_ALIGN.LEFT)
    txt(s, Inches(0.4), Inches(7.08), Inches(7.5), Inches(0.30),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}", 9, GREY, False)]])
    txt(s, Inches(5.0), Inches(7.08), Inches(3.3), Inches(0.30),
        [[("© 2026 Tertiary Infotech Pte Ltd", 9, GREY, False)]], align=PP_ALIGN.CENTER)
    txt(s, Inches(12.4), Inches(7.08), Inches(0.6), Inches(0.30),
        [[(str(PAGE["n"]), 9, GREY, False)]], align=PP_ALIGN.RIGHT)


def head(s, title, kicker=None, kcolor=BLUE, title_w=11.9, title_size=28):
    rect(s, 0, 0, SW, SH, WHITE); rect(s, 0, 0, Inches(0.28), Inches(1.55), kcolor)
    if kicker:
        txt(s, Inches(0.85), Inches(0.34), Inches(11.6), Inches(0.38), [[(kicker, 14, kcolor, True)]])
    txt(s, Inches(0.85), Inches(0.78), Inches(title_w), Inches(0.86),
        [[(title, title_size, INK, True)]])
    rect(s, Inches(0.85), Inches(1.68), Inches(11.63), Inches(0.02), LINE)
    return s


def _asset(name):
    p = os.path.join(ASSETS, name)
    return p if os.path.exists(p) else None


def _fit(path, max_w_in, max_h_in):
    """Aspect-fit an image to a box, returning (w,h) in EMU. Never stretches."""
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    w = max_w_in; h = w / ar
    if h > max_h_in:
        h = max_h_in; w = h * ar
    return Inches(w), Inches(h)


def _picture_cover(s, path, x, y, w, h):
    """Fill a fixed picture frame without stretching the source image."""
    with Image.open(path) as im:
        iw, ih = im.size
    image_ar = iw / ih
    frame_ar = w / h
    pic = s.shapes.add_picture(path, x, y)
    pic.width = int(w); pic.height = int(h)
    if image_ar > frame_ar:
        crop = (1 - frame_ar / image_ar) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - image_ar / frame_ar) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic


# ================================================================ components
def cover():
    s = slide(); rect(s, 0, 0, SW, SH, NAVY)
    hero = _asset("hero-ai-agent-security-v30.png") if C.VERSION.startswith("3") else _asset("hero-ai-agent-security-v21.png")
    if hero:
        s.shapes.add_picture(hero, 0, 0, width=SW, height=SH)
    rect(s, 0, 0, Inches(6.45), SH, NAVY)
    rect(s, 0, 0, SW, Inches(0.16), BLUE); rect(s, 0, Inches(7.34), SW, Inches(0.16), RED)
    org = _asset("tertiary-infotech-logo.png")
    if org:
        # The source PNG already has a transparent alpha channel. Keep the mark directly
        # on the navy cover instead of placing it on a white backing rectangle.
        s.shapes.add_picture(org, Inches(0.82), Inches(0.62), height=Inches(0.78))
    rect(s, Inches(10.72), Inches(0.62), Inches(1.85), Inches(0.94), RED)
    txt(s, Inches(10.72), Inches(0.73), Inches(1.85), Inches(0.46),
        [[("WSQ", 22, WHITE, True)]], align=PP_ALIGN.CENTER)
    txt(s, Inches(10.72), Inches(1.18), Inches(1.85), Inches(0.32),
        [[("AI SECURITY", 8, WHITE, True)]], align=PP_ALIGN.CENTER)
    txt(s, Inches(0.8), Inches(1.9), Inches(5.2), Inches(0.5),
        [[("TRAINER SLIDES  ·  WSQ", 15, CODEBLUE, True)]])
    txt(s, Inches(0.8), Inches(2.35), Inches(5.35), Inches(2.05), [[(C.TITLE, 34, WHITE, True)]])
    rect(s, Inches(0.82), Inches(4.35), Inches(2.2), Inches(0.06), RED)
    txt(s, Inches(0.8), Inches(4.62), Inches(5.35), Inches(1.35),
        [[(f"WSQ Course Code: {C.COURSE_CODE}  ·  {C.DURATION}", 14, WHITE, False)],
         [(f"Skills Framework TSC: {C.TSC_TITLE} ({C.TSC_CODE})", 12, WHITE, False)],
         [("Conducted by Tertiary Infotech Pte Ltd  ·  UEN 201200696W", 12, WHITE, False)]], space=6)
    txt(s, Inches(0.8), Inches(6.2), Inches(5.35), Inches(0.36),
        [[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}", 11, CODEBLUE, False)]])
    txt(s, Inches(0.8), Inches(6.68), Inches(5.35), Inches(0.28),
        [[("© 2026 Tertiary Infotech Pte Ltd. All Rights Reserved.  ·  www.tertiarycourses.com.sg",
           8.5, WHITE, False)]])
    PAGE["n"] += 1


def section(kicker, title, n, sub=""):
    s = slide(); rect(s, 0, 0, SW, SH, WHITE); rect(s, 0, 0, Inches(0.28), SH, BLUE)
    rect(s, Inches(0.85), Inches(2.5), Inches(0.14), Inches(2.0), RED)
    txt(s, Inches(1.25), Inches(2.55), Inches(11), Inches(0.6), [[(kicker, 18, BLUE, True)]])
    txt(s, Inches(1.25), Inches(3.0), Inches(11.4), Inches(1.6), [[(title, 40, INK, True)]])
    if sub: txt(s, Inches(1.27), Inches(4.55), Inches(11), Inches(0.8), [[(sub, 16, GREY, False)]])
    txt(s, Inches(10.0), Inches(0.7), Inches(2.8), Inches(1.6),
        [[(n, 72, LINE, True)]], align=PP_ALIGN.RIGHT)
    footer(s)


def content(title, items, kicker=None, size=20, kcolor=BLUE):
    s = head(slide(), title, kicker, kcolor=kcolor)
    bullets(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(4.9), items, size=size)
    footer(s); return s


def two_col(title, left, right, kicker=None, lhead="", rhead="",
            lcolor=BLUE, rcolor=TEAL, note=None):
    s = head(slide(), title, kicker)
    bh = Inches(4.7) if not note else Inches(4.05)
    rect(s, Inches(0.85), Inches(1.95), Inches(5.7), bh, LIGHT)
    rect(s, Inches(6.95), Inches(1.95), Inches(5.55), bh, LIGHT)
    rect(s, Inches(0.85), Inches(1.95), Inches(5.7), Inches(0.1), lcolor)
    rect(s, Inches(6.95), Inches(1.95), Inches(5.55), Inches(0.1), rcolor)
    if lhead: txt(s, Inches(1.1), Inches(2.15), Inches(5.2), Inches(0.4), [[(lhead, 16, lcolor, True)]])
    if rhead: txt(s, Inches(7.2), Inches(2.15), Inches(5.0), Inches(0.4), [[(rhead, 16, rcolor, True)]])
    bullets(s, Inches(1.1), Inches(2.7), Inches(5.2), bh - Inches(0.9), left, size=15)
    bullets(s, Inches(7.2), Inches(2.7), Inches(5.05), bh - Inches(0.9), right, size=15, mcolor=rcolor)
    if note:
        txt(s, Inches(0.85), Inches(6.14), Inches(11.7), Inches(0.50), [[(note, 12.3, GREY, False)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def cards3(title, cards, kicker, kcolor=BLUE):
    s = head(slide(), title, kicker, kcolor=kcolor); xs = [Inches(0.85), Inches(5.0), Inches(9.15)]
    for i, c in enumerate(cards[:3]):
        x = xs[i]; col = c[0]
        rect(s, x, Inches(1.95), Inches(3.65), Inches(4.7), LIGHT)
        rect(s, x, Inches(1.95), Inches(3.65), Inches(0.12), col)
        txt(s, x + Inches(0.25), Inches(2.2), Inches(3.2), Inches(0.6), [[(c[1], 19, col, True)]])
        bullets(s, x + Inches(0.25), Inches(2.95), Inches(3.2), Inches(3.4), c[2],
                size=14, mcolor=col, gap=9)
    footer(s); return s


def big_statement(line1, line2, kicker, color=BLUE):
    s = slide(); rect(s, 0, 0, SW, SH, WHITE); rect(s, 0, 0, Inches(0.28), SH, color)
    txt(s, Inches(1.1), Inches(2.2), Inches(11), Inches(0.5), [[(kicker, 16, color, True)]])
    txt(s, Inches(1.1), Inches(2.8), Inches(11.3), Inches(2.4), [[(line1, 38, INK, True)]])
    if line2: txt(s, Inches(1.12), Inches(4.9), Inches(11), Inches(1.2), [[(line2, 20, GREY, False)]])
    footer(s); return s


def tile_grid(title, items, kicker=None, cols=2, size=15, accent=BLUE):
    s = head(slide(), title, kicker, kcolor=accent)
    n = len(items); rows = math.ceil(n / cols)
    X0 = Inches(0.85); Y0 = Inches(1.95); TOTW = Inches(11.63); AREAH = Inches(4.78)
    gx = Inches(0.25); gy = Inches(0.22)
    cw = int((TOTW - gx * (cols - 1)) / cols); ch = int((AREAH - gy * (rows - 1)) / rows)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        x = int(X0 + (cw + gx) * c); y = int(Y0 + (ch + gy) * r)
        col = PALETTE[i % len(PALETTE)] if accent is None else accent
        col = PALETTE[i % len(PALETTE)]
        rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, Inches(0.09), ch, col)
        if isinstance(it, tuple):
            txt(s, x + Inches(0.3), y + Inches(0.16), cw - Inches(0.55), ch - Inches(0.3),
                [[(it[0], size, col, True)], [(it[1], size - 3, INK, False)]], space=3)
        else:
            txt(s, x + Inches(0.3), y, cw - Inches(0.55), ch,
                [[(it, size, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def flow_h(title, steps, kicker=None, color=BLUE, note=None):
    s = head(slide(), title, kicker, kcolor=color)
    n = len(steps); X0 = Inches(0.85); TOTW = Inches(11.63); gap = Inches(0.34)
    cw = int((TOTW - gap * (n - 1)) / n); y = Inches(2.45); ch = Inches(3.15); bd = Inches(0.82)
    for i, st in enumerate(steps):
        x = int(X0 + (cw + gap) * i)
        rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, cw, Inches(0.1), color)
        oval(s, int(x + cw / 2 - bd / 2), int(y + Inches(0.42)), bd, bd, color)
        txt(s, int(x + cw / 2 - bd / 2), int(y + Inches(0.42)), bd, bd,
            [[(str(i + 1), 30, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.16), int(y + Inches(1.55)), cw - Inches(0.32), int(ch - Inches(1.7)),
            [[(st, 13, INK, False)]], align=PP_ALIGN.CENTER)
        if i < n - 1:
            txt(s, int(x + cw - Inches(0.04)), int(y + ch / 2 - Inches(0.3)),
                int(gap + Inches(0.08)), Inches(0.6),
                [[("▶", 15, color, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if note:
        txt(s, Inches(0.85), Inches(5.95), Inches(11.7), Inches(0.7), [[(note, 14, GREY, False)]],
            align=PP_ALIGN.CENTER)
    footer(s); return s


def trainer_slide(kicker, name, role, rows, initials, accent=BLUE):
    s = head(slide(), "About the Trainer", kicker, kcolor=accent)
    lx = Inches(0.85); lw = Inches(3.65)
    rect(s, lx, Inches(1.95), lw, Inches(4.7), LIGHT); rect(s, lx, Inches(1.95), lw, Inches(0.12), accent)
    bd = Inches(1.7); ax = int(lx + (lw - bd) / 2)
    oval(s, ax, Inches(2.5), bd, bd, accent)
    txt(s, ax, Inches(2.5), bd, bd, [[(initials, 44, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx + Inches(0.15), Inches(4.55), lw - Inches(0.3), Inches(0.6),
        [[(name, 21, INK, True)]], align=PP_ALIGN.CENTER)
    txt(s, lx + Inches(0.15), Inches(5.2), lw - Inches(0.3), Inches(1.2),
        [[(role, 13, GREY, False)]], align=PP_ALIGN.CENTER)
    rx = Inches(4.9); rw = Inches(7.6); ry = Inches(1.95); rh = Inches(4.7)
    n = len(rows); gy = Inches(0.2); th = int((rh - gy * (n - 1)) / n)
    for i, (label, val) in enumerate(rows):
        y = int(ry + (th + gy) * i); col = PALETTE[i % len(PALETTE)]
        rect(s, rx, y, rw, th, LIGHT); rect(s, rx, y, Inches(0.1), th, col)
        vruns = [(val, 14, INK, False)] if val else \
                [("____________________________________________", 13, LINE, False)]
        txt(s, rx + Inches(0.32), y, rw - Inches(0.6), th,
            [[(label.upper(), 11, col, True)], vruns], anchor=MSO_ANCHOR.MIDDLE, space=3)
    footer(s); return s


def brk(kind, dur, color=AMBER):
    s = slide(); rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, SW, Inches(0.22), color)
    rect(s, Inches(5.4), Inches(2.35), Inches(2.53), Inches(0.1), color)
    txt(s, 0, Inches(2.75), SW, Inches(1.2), [[(kind, 48, INK, True)]], align=PP_ALIGN.CENTER)
    txt(s, 0, Inches(4.05), SW, Inches(0.8), [[(dur, 22, color, True)]], align=PP_ALIGN.CENTER)
    # Break slides carry the same footer as every other slide (house rule: every slide).
    footer(s)
    return s


def img_points(title, image, points, kicker=None, accent=BLUE, img_w=7.0, note=None,
               frame="dark"):
    s = head(slide(), title, kicker, kcolor=accent)
    p = _asset(image)
    content_h = 4.08 if note else 4.75
    if p:
        ix = Inches(0.85); iy = Inches(2.05); iw = Inches(img_w); ih = Inches(content_h)
        frame_pad = Inches(0.03) if frame == "soft" else Inches(0.05)
        frame_color = LINE if frame == "soft" else NAVY
        rect(s, ix - frame_pad, iy - frame_pad,
             iw + frame_pad * 2, ih + frame_pad * 2, frame_color)
        _picture_cover(s, p, ix, iy, iw, ih)
    rx = Inches(0.85) + Inches(img_w) + Inches(0.3); rw = Inches(12.48) - rx
    n = len(points); gy = Inches(0.2); th = int((Inches(content_h) - gy * (n - 1)) / n)
    for i, (t1, t2) in enumerate(points):
        y = int(Inches(2.05) + (th + gy) * i); col = PALETTE[i % len(PALETTE)]
        rect(s, rx, y, rw, th, LIGHT); rect(s, rx, y, Inches(0.09), th, col)
        txt(s, rx + Inches(0.28), y, rw - Inches(0.5), th,
            [[(t1, 14, col, True)], [(t2, 12, INK, False)]], anchor=MSO_ANCHOR.MIDDLE, space=3)
    if note:
        txt(s, Inches(0.85), Inches(6.20), Inches(11.7), Inches(0.44), [[(note, 9.8, GREY, False)]],
            align=PP_ALIGN.CENTER)
    footer(s); return s


def fashion_case_slide(title, image, cards, kicker=None, accent=TEAL):
    """Editorial portrait plus three evidence cards for the AI-fashion case."""
    s = head(slide(), title, kicker, kcolor=accent)
    ix = Inches(0.85); iy = Inches(1.96); iw = Inches(4.42); ih = Inches(4.66)
    rect(s, ix - Inches(0.05), iy - Inches(0.05), iw + Inches(0.10), ih + Inches(0.10), NAVY)
    p = _asset(image)
    if p:
        pic = _picture_cover(s, p, ix, iy, iw, ih)
        # Bias the portrait crop upward so the model's face remains fully visible.
        total_crop = pic.crop_top + pic.crop_bottom
        pic.crop_top = min(0.035, total_crop)
        pic.crop_bottom = max(0, total_crop - pic.crop_top)
    rect(s, ix, iy + ih - Inches(0.52), iw, Inches(0.52), NAVY)
    txt(s, ix + Inches(0.18), iy + ih - Inches(0.42), iw - Inches(0.36), Inches(0.28),
        [[("SYNTHETIC MODEL · IMAGEGEN", 9.5, WHITE, True)]], align=PP_ALIGN.CENTER)
    rx = Inches(5.62); rw = Inches(6.86); gap = Inches(0.18); ch = Inches(1.43)
    for i, (label, body) in enumerate(cards):
        y = iy + i * (ch + gap); col = PALETTE[i % len(PALETTE)]
        rect(s, rx, y, rw, ch, LIGHT); rect(s, rx, y, Inches(0.10), ch, col)
        oval(s, rx + Inches(0.22), y + Inches(0.22), Inches(0.48), Inches(0.48), col)
        txt(s, rx + Inches(0.22), y + Inches(0.22), Inches(0.48), Inches(0.48),
            [[(str(i + 1), 12, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, rx + Inches(0.86), y + Inches(0.16), rw - Inches(1.10), Inches(0.34),
            [[(label, 14.5, col, True)]])
        txt(s, rx + Inches(0.86), y + Inches(0.55), rw - Inches(1.10), Inches(0.72),
            [[(body, 11.8, INK, False)]])
    footer(s); return s


def img_full(title, image, kicker=None, accent=BLUE, caption=None):
    s = head(slide(), title, kicker, kcolor=accent)
    p = _asset(image)
    maxh = 4.35 if caption else 4.85
    if p:
        w, h = _fit(p, 11.6, maxh)
        x = int(Inches(0.85) + (Inches(11.63) - w) / 2)
        s.shapes.add_picture(p, x, Inches(1.95), width=w, height=h)
    if caption:
        rect(s, Inches(0.85), Inches(6.16), Inches(11.63), Inches(0.48), LIGHT)
        txt(s, Inches(1.1), Inches(6.16), Inches(11.1), Inches(0.48), [[(caption, 11.8, INK, False)]],
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    footer(s); return s


def _embed_youtube(s, poster_path, embed_url, watch_url, x, y, w, h):
    """Insert a real PowerPoint 'Online Video' (YouTube) object so it streams and
    plays in-slide during the slideshow — the same structure PowerPoint writes when
    you use Insert > Video > Online Video. The poster image is the video frame."""
    from pptx.oxml.ns import qn
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    part = s.part
    # 1) external relationship to the YouTube embed URL (the streamed media)
    media_rid = part.relate_to(embed_url, RT.MEDIA, is_external=True)
    # 2) a second external relationship PowerPoint also stores (video web link)
    video_rid = part.relate_to(watch_url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video", is_external=True)
    # 3) the poster image as an internal picture part
    img_part, image_rid = part.get_or_add_image_part(poster_path)
    ns = ('xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
          'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
          'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"')
    xml = (
        f'<p:pic {ns}>'
        f'<p:nvPicPr>'
        f'<p:cNvPr id="{s.shapes._next_shape_id}" name="YouTube Online Video">'
        f'<a:hlinkClick r:id="" action="ppaction://media"/>'
        f'</p:cNvPr>'
        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        f'<p:nvPr>'
        f'<a:videoFile r:link="{video_rid}"/>'
        f'<p:extLst>'
        f'<p:ext uri="{{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}}">'
        f'<p14:media xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" r:embed="{media_rid}"/>'
        f'</p:ext>'
        f'</p:extLst>'
        f'</p:nvPr>'
        f'</p:nvPicPr>'
        f'<p:blipFill><a:blip r:embed="{image_rid}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{int(x)}" y="{int(y)}"/><a:ext cx="{int(w)}" cy="{int(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'</p:spPr>'
        f'</p:pic>'
    )
    from pptx.oxml import parse_xml
    pic = parse_xml(xml)
    s.shapes._spTree.append(pic)
    return pic


def video_slide(title, poster, video_url, kicker=None, accent=BLUE, caption=None,
                points=None, video_embed=None, video_file=None):
    """A slide with an embedded video that plays in-slide in PowerPoint slideshow.
    If video_file is given it is embedded as a real local movie (plays offline);
    otherwise a YouTube online-video object is used. A clickable link is the fallback
    for PDF export and non-PowerPoint viewers. The vertical clip sits in a portrait
    frame on the left with talking points on the right."""
    s = head(slide(), title, kicker, kcolor=accent)
    # Size the video frame from the poster's aspect ratio so both portrait (9:16)
    # and landscape (16:9) clips sit correctly. Portrait fills the height; landscape
    # fills a wider box and drops the frame lower so the title/banner stay clear.
    poster_probe = _asset(poster) if poster else None
    ar = 9 / 16
    if poster_probe:
        try:
            with Image.open(poster_probe) as _im:
                ar = _im.size[0] / _im.size[1]
        except Exception:
            pass
    if ar >= 1:  # landscape
        vw = Inches(4.65); vh = Inches(vw / 914400 * (1 / ar)); vx = Inches(1.05); vy = Inches(2.35)
    else:        # portrait
        vh = Inches(4.30); vw = Inches(vh / 914400 * ar); vx = Inches(1.22); vy = Inches(2.02)
    rect(s, int(vx - Inches(0.1)), int(vy - Inches(0.1)),
         int(vw + Inches(0.2)), int(vh + Inches(0.2)), NAVY)
    local = bool(video_file and _asset(video_file))
    banner = "EMBEDDED VIDEO · PLAYS IN SLIDESHOW" if local else "ONLINE VIDEO · YOUTUBE"
    rect(s, int(vx - Inches(0.1)), int(vy - Inches(0.36)), int(vw + Inches(0.2)), Inches(0.30), RED)
    txt(s, int(vx - Inches(0.1)), int(vy - Inches(0.34)), int(vw + Inches(0.2)), Inches(0.24),
        [[(banner, 8.0, WHITE, True)]], align=PP_ALIGN.CENTER)
    poster_path = _asset(poster) if poster else None
    embedded = False
    if local:
        try:
            s.shapes.add_movie(_asset(video_file), vx, vy, vw, vh,
                               poster_frame_image=poster_path, mime_type="video/mp4")
            embedded = True
        except Exception as e:
            print(f"  [video] local movie embed failed ({e}); trying poster fallback")
    if not embedded and poster_path and video_embed:
        try:
            _embed_youtube(s, poster_path, video_embed, video_url, vx, vy, vw, vh)
            embedded = True
        except Exception as e:
            print(f"  [video] online embed failed ({e}); using poster fallback")
    if not embedded and poster_path:
        s.shapes.add_picture(poster_path, vx, vy, width=vw, height=vh)
    if not embedded:
        bd = Inches(0.9)
        oval(s, int(vx + vw/2 - bd/2), int(vy + vh/2 - bd/2), bd, bd, RED)
        txt(s, int(vx + vw/2 - bd/2), int(vy + vh/2 - bd/2), bd, bd,
            [[("▶", 26, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # right-hand talking points — a fixed band so they read full for both portrait
    # and landscape video frames.
    rx = int(vx + vw + Inches(0.55)); rw = int(Inches(12.48) - rx)
    if points:
        band_y = Inches(1.95); band_h = Inches(4.55)
        n = len(points); gy = Inches(0.2); th = int((band_h - gy*(n-1))/n)
        for i, (t1, t2) in enumerate(points):
            y = int(band_y + (th+gy)*i); col = PALETTE[i % len(PALETTE)]
            rect(s, rx, y, rw, th, LIGHT); rect(s, rx, y, Inches(0.09), th, col)
            txt(s, rx + Inches(0.28), y, rw - Inches(0.5), th,
                [[(t1, 15, col, True)], [(t2, 12.5, INK, False)]],
                anchor=MSO_ANCHOR.MIDDLE, space=3)
    # link line under the video: for a local movie there is no external link to show
    if video_url and not local:
        link_tb = txt(s, vx - Inches(0.1), int(vy + vh + Inches(0.12)), int(vw + Inches(0.2)),
                      Inches(0.34), [[("▶ Open on YouTube (fallback)", 10.5, BLUE, True)]],
                      align=PP_ALIGN.CENTER)
        try:
            link_tb.text_frame.paragraphs[0].runs[0].hyperlink.address = video_url
        except Exception:
            pass
    elif local:
        txt(s, vx - Inches(0.1), int(vy + vh + Inches(0.12)), int(vw + Inches(0.2)),
            Inches(0.34), [[("Click the video, then Play, in slideshow", 9.5, GREY, False)]],
            align=PP_ALIGN.CENTER)
    if caption:
        txt(s, rx, Inches(6.42), rw, Inches(0.28),
            [[(caption, 9.4, GREY, False)]], align=PP_ALIGN.CENTER)
    footer(s); return s


def table_slide(title, headers, rows, kicker=None, accent=BLUE, widths=None, note=None, fsize=13):
    s = head(slide(), title, kicker, kcolor=accent)
    ncol = len(headers); X0 = Inches(0.85); TOTW = Inches(11.63)
    ws = [int(TOTW * w) for w in widths] if widths else [int(TOTW / ncol)] * ncol
    area_h = Inches(4.08) if note else Inches(4.85)
    nrow = len(rows) + 1; rh = int(area_h / nrow); y = Inches(1.95); x = X0
    for j, htxt in enumerate(headers):
        rect(s, x, y, ws[j], rh, accent)
        txt(s, x + Inches(0.14), y, ws[j] - Inches(0.24), rh, [[(htxt, fsize, WHITE, True)]],
            anchor=MSO_ANCHOR.MIDDLE)
        x += ws[j]
    for i, row in enumerate(rows):
        y = int(Inches(1.95) + rh * (i + 1)); x = X0
        fill = LIGHT if i % 2 == 0 else WHITE
        for j, cell in enumerate(row):
            rect(s, x, y, ws[j], rh, fill, line=LINE)
            txt(s, x + Inches(0.14), y, ws[j] - Inches(0.24), rh,
                [[(cell, fsize - 1, INK, j == 0)]], anchor=MSO_ANCHOR.MIDDLE)
            x += ws[j]
    if note:
        txt(s, Inches(0.85), Inches(6.16), Inches(11.7), Inches(0.48), [[(note, 11.8, GREY, False)]],
            align=PP_ALIGN.CENTER)
    footer(s); return s


def formula_slide(title, panels, kicker=None, accent=BLUE, note=None):
    s = head(slide(), title, kicker, kcolor=accent)
    n = len(panels); X0 = Inches(0.85); TOTW = Inches(11.63); gap = Inches(0.3)
    cw = int((TOTW - gap * (n - 1)) / n); y = Inches(2.1); ch = Inches(3.9)
    for i, (hd, formula, cap) in enumerate(panels):
        x = int(X0 + (cw + gap) * i); col = PALETTE[i % len(PALETTE)]
        rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, cw, Inches(0.1), col)
        txt(s, x + Inches(0.22), y + Inches(0.28), cw - Inches(0.44), Inches(0.5),
            [[(hd, 16, col, True)]])
        rect(s, x + Inches(0.22), y + Inches(0.95), cw - Inches(0.44), Inches(1.35), NAVY)
        txt(s, x + Inches(0.32), y + Inches(0.95), cw - Inches(0.64), Inches(1.35),
            [[(ln, 15, CODEBLUE, True)] for ln in formula.split("\n")],
            anchor=MSO_ANCHOR.MIDDLE, space=3)
        txt(s, x + Inches(0.22), y + Inches(2.5), cw - Inches(0.44), ch - Inches(2.6),
            [[(cap, 12.5, GREY, False)]])
    if note:
        txt(s, Inches(0.85), Inches(6.14), Inches(11.7), Inches(0.50), [[(note, 12.3, GREY, False)]],
            align=PP_ALIGN.CENTER)
    footer(s); return s


def ncards(title, cards, kicker=None, accent=BLUE, cols=4, note=None, ch_in=None):
    """Outlined numbered concept cards — the reference's signature 4-across move.

    The title sits to the right of the badge and may wrap to two lines; the body
    starts BELOW whichever is taller, so a long title can never overlap the body.
    """
    s = head(slide(), title, kicker, kcolor=accent)
    n = len(cards); X0 = Inches(0.72); TOTW = Inches(11.85)
    gx = Inches(0.28) if cols > 2 else Inches(0.15)
    cw = int((TOTW - gx * (cols - 1)) / cols)
    rows = math.ceil(n / cols)
    areah = Inches(3.95) if note else Inches(4.75)
    gy = Inches(0.25)
    ch = Inches(ch_in) if ch_in else int((areah - gy * (rows - 1)) / rows)
    # card title column width, in characters, drives the wrap estimate
    title_w_in = (cw - Inches(1.0)) / 914400
    for i, (t1, t2) in enumerate(cards):
        r, c = divmod(i, cols)
        x = int(X0 + (cw + gx) * c); y = int(Inches(1.95) + (ch + gy) * r)
        col = PALETTE[i % len(PALETTE)]
        rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, Inches(0.09), ch, col)
        bd = Inches(0.46)
        oval(s, x + Inches(0.24), y + Inches(0.20), bd, bd, col)
        txt(s, x + Inches(0.24), y + Inches(0.20), bd, bd, [[(str(i + 1), 15, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # ~ 9.5 chars per inch at 14pt Arial bold
        est_lines = max(1, math.ceil(len(t1) / max(6, title_w_in * 9.5)))
        thh = Inches(0.30) * est_lines
        txt(s, x + Inches(0.82), y + Inches(0.16), cw - Inches(1.0), thh,
            [[(t1, 14, col, True)]])
        body_y = y + Inches(0.16) + max(thh, Inches(0.50)) + Inches(0.16)
        txt(s, x + Inches(0.28), body_y, cw - Inches(0.54),
            y + ch - body_y - Inches(0.14), [[(t2, 11.5, INK, False)]])
    if note:
        rect(s, Inches(0.72), Inches(6.10), Inches(11.85), Inches(0.54), LIGHT)
        txt(s, Inches(1.0), Inches(6.10), Inches(11.3), Inches(0.54), [[(note, 11.7, INK, False)]],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def takeaway(s, text, color=BLUE, y=6.35):
    """One bolded boxed rule to remember — closes a concept slide."""
    rect(s, Inches(0.85), Inches(y), Inches(11.63), Inches(0.62), LIGHT)
    rect(s, Inches(0.85), Inches(y), Inches(0.09), Inches(0.62), color)
    txt(s, Inches(1.15), Inches(y), Inches(11.2), Inches(0.62), [[(text, 13.5, INK, True)]],
        anchor=MSO_ANCHOR.MIDDLE)


def activity_slide(a):
    col = CMAP[a["accent"]]
    # Reserve the top-right area for the duration badge. Long activity titles must
    # never run underneath it.
    s = head(slide(), a["title"], f"ACTIVITY {a['n']} · HANDS-ON", kcolor=col,
             title_w=9.15, title_size=26)
    rect(s, Inches(10.35), Inches(0.5), Inches(2.13), Inches(0.62), col)
    txt(s, Inches(10.35), Inches(0.5), Inches(2.13), Inches(0.62),
        [[(f"ACTIVITY {a['n']} · {a['minutes']} MIN", 11, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # org + scenario
    rect(s, Inches(0.85), Inches(1.95), Inches(11.63), Inches(1.32), LIGHT)
    rect(s, Inches(0.85), Inches(1.95), Inches(0.09), Inches(1.32), col)
    txt(s, Inches(1.15), Inches(2.06), Inches(11.2), Inches(0.4), [[(a["org"], 13, col, True)]])
    txt(s, Inches(1.15), Inches(2.44), Inches(11.2), Inches(0.8),
        [[(a["scenario"], 13.5, INK, False)]])
    # 5 chips
    n = len(a["flow"]); X0 = Inches(0.85); TOTW = Inches(11.63); gap = Inches(0.22)
    cw = int((TOTW - gap * (n - 1)) / n); y = Inches(3.5); ch = Inches(1.25); bd = Inches(0.56)
    for i, st in enumerate(a["flow"]):
        x = int(X0 + (cw + gap) * i)
        rect(s, x, y, cw, ch, WHITE, line=LINE)
        oval(s, int(x + cw / 2 - bd / 2), y + Inches(0.14), bd, bd, col)
        txt(s, int(x + cw / 2 - bd / 2), y + Inches(0.14), bd, bd, [[(str(i + 1), 15, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.1), y + Inches(0.76), cw - Inches(0.2), Inches(0.45),
            [[(st, 10.5, INK, False)]], align=PP_ALIGN.CENTER)
        if i < n - 1:
            txt(s, int(x + cw), int(y + ch / 2 - Inches(0.22)), int(gap), Inches(0.44),
                [[("▶", 12, col, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # discussion questions
    rect(s, Inches(0.85), Inches(4.95), Inches(7.6), Inches(1.42), WHITE, line=LINE)
    txt(s, Inches(1.1), Inches(5.02), Inches(7.2), Inches(0.32),
        [[("DISCUSSION QUESTIONS", 10.5, col, True)]])
    bullets(s, Inches(1.1), Inches(5.32), Inches(7.15), Inches(1.0), a["questions"],
            size=10.5, gap=3)
    # produce band
    rect(s, Inches(8.65), Inches(4.95), Inches(3.83), Inches(1.42), LIGHT)
    rect(s, Inches(8.65), Inches(4.95), Inches(0.09), Inches(1.42), TEAL)
    txt(s, Inches(8.92), Inches(5.02), Inches(3.4), Inches(0.32),
        [[("YOU'LL PRODUCE", 10.5, TEAL, True)]])
    txt(s, Inches(8.92), Inches(5.32), Inches(3.45), Inches(1.0), [[(a["produce"], 11.5, INK, False)]])
    # meta
    txt(s, Inches(0.85), Inches(6.30), Inches(11.63), Inches(0.34),
        [[(f"Assesses {a['ka']}  ·  Activity pack: activities/{a['folder']}/  "
           f"·  Full step-by-step facilitation detail: Learner Guide", 11.5, GREY, False)]])
    footer(s); return s


def lms_slide(spec=None):
    spec = spec or {}
    s = head(slide(), spec.get("title", "Download Course Material"),
             kicker="COURSE PORTAL · LMS/TMS", kcolor=BLUE)
    bx, by, bw, bh = Inches(0.85), Inches(2.0), Inches(6.1), Inches(4.55)
    rect(s, bx, by, bw, bh, WHITE, line=LINE)
    rect(s, bx, by, bw, Inches(0.52), RGBColor(0xEE, 0xF2, 0xF8))
    for i, c in enumerate([RED, AMBER, TEAL]):
        oval(s, int(bx + Inches(0.18) + Inches(0.28) * i), int(by + Inches(0.17)),
             Inches(0.18), Inches(0.18), c)
    rect(s, int(bx + Inches(1.2)), int(by + Inches(0.1)), int(bw - Inches(1.5)), Inches(0.34),
         WHITE, line=LINE)
    txt(s, int(bx + Inches(1.35)), int(by + Inches(0.1)), int(bw - Inches(1.8)), Inches(0.34),
        [[("https://lms-tms.tertiaryinfotech.com", 13, BLUE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    shot = _asset("lms-tms-login.png")
    if shot:
        w, h = _fit(shot, 5.6, 3.55)
        x = int(bx + (bw - w) / 2)
        s.shapes.add_picture(shot, x, int(by + Inches(0.72)), width=w, height=h)
    supplied = spec.get("points", [])
    if supplied:
        steps = [(f"Step {i}", point) for i, point in enumerate(supplied, 1)]
    else:
        steps = [("Sign in", "lms-tms.tertiaryinfotech.com — log in with your registered email (OTP or password)."),
                 ("Open your course", f"Select '{C.TITLE}' under My Courses."),
                 ("Download materials", "Trainer slides and the Learner Guide — your open-book references."),
                 ("Submit & survey", "Upload your assessment answers and complete the TRAQOM survey.")]
    rx = Inches(7.3); rw = Inches(5.2); gy = Inches(0.2)
    th = int((Inches(4.55) - gy * (len(steps) - 1)) / len(steps))
    for i, (t1, t2) in enumerate(steps):
        y = int(Inches(2.0) + (th + gy) * i); col = PALETTE[i % 4]
        rect(s, rx, y, rw, th, LIGHT); rect(s, rx, y, Inches(0.09), th, col)
        bd = Inches(0.5)
        oval(s, rx + Inches(0.2), int(y + th / 2 - bd / 2), bd, bd, col)
        txt(s, rx + Inches(0.2), int(y + th / 2 - bd / 2), bd, bd, [[(str(i + 1), 16, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, rx + Inches(0.9), y, rw - Inches(1.1), th,
            [[(t1, 14, col, True)], [(t2, 11.5, INK, False)]], anchor=MSO_ANCHOR.MIDDLE, space=2)
    note = spec.get("note") or ("Use only the current learner-facing files. Keep evidence "
                                "synthetic, local and reversible; never upload secrets or answer keys.")
    txt(s, Inches(0.85), Inches(6.20), Inches(11.7), Inches(0.44),
        [[(note, 12.5, GREY, False)]], align=PP_ALIGN.CENTER)
    footer(s); return s


def attendance_slide(kicker="TRAQOM · SSG DIGITAL ATTENDANCE"):
    return content("Digital Attendance (Mandatory)", [
        "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
        "The trainer or administrator displays the digital attendance QR code from the SSG portal.",
        "Scan the QR code with your mobile phone camera and submit your attendance.",
        "A minimum of 75% attendance is required to be eligible for assessment and funding.",
        "Complete the TRAQOM survey at the end of the course — it is required for funding.",
    ], kicker=kicker)


def _v30_thank_you(spec):
    s = slide(); rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, SW, Inches(0.22), BLUE); rect(s, 0, Inches(7.28), SW, Inches(0.22), RED)
    txt(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.25),
        [[(spec.get("title", "Thank You"), 46, INK, True)]], align=PP_ALIGN.CENTER)
    rect(s, Inches(5.4), Inches(3.85), Inches(2.53), Inches(0.08), RED)
    txt(s, Inches(1.15), Inches(4.25), Inches(11.0), Inches(1.15),
        [[(spec.get("note", "Bound the authority. Preserve the evidence."), 21, GREY, False)]],
        align=PP_ALIGN.CENTER)
    txt(s, 0, Inches(5.55), SW, Inches(0.6),
        [[(f"{C.COURSE_CODE}  ·  Tertiary Infotech Pte Ltd", 13, GREY, False)]],
        align=PP_ALIGN.CENTER)
    footer(s)


def _v30_render(spec):
    """Render one v3 source record with the existing house component library."""
    kind = spec.get("kind", "content")
    title = spec.get("title", "")
    kicker = spec.get("kicker")
    accent = CMAP.get(spec.get("accent", "BLUE"), BLUE)
    note = spec.get("note")

    if kind == "cover":
        cover()
    elif kind == "section":
        section(kicker or "COURSE MODULE", title,
                spec.get("number", str(spec.get("n", "")).zfill(2)), note or "")
    elif kind == "content":
        content(title, spec.get("points", []), kicker=kicker, size=spec.get("size", 17), kcolor=accent)
    elif kind == "cards":
        cards = []
        for item in spec.get("cards", []):
            if isinstance(item, dict):
                cards.append((item.get("title", ""), item.get("body", "")))
            else:
                cards.append(tuple(item[:2]))
        ncards(title, cards, kicker=kicker, accent=accent,
               cols=spec.get("cols", min(4, max(2, len(cards)))), note=note,
               ch_in=spec.get("card_height"))
    elif kind == "fashion_case":
        cards = []
        for item in spec.get("cards", []):
            cards.append(tuple(item[:2]) if not isinstance(item, dict)
                         else (item.get("title", ""), item.get("body", "")))
        fashion_case_slide(title, spec.get("image"), cards, kicker=kicker, accent=accent)
    elif kind == "table":
        table_slide(title, spec.get("headers", []), spec.get("rows", []),
                    kicker=kicker, accent=accent, widths=spec.get("widths"),
                    note=note, fsize=spec.get("fsize", 12.5))
    elif kind == "compare":
        two_col(title, spec.get("left", []), spec.get("right", []), kicker=kicker,
                lhead=spec.get("lhead", ""), rhead=spec.get("rhead", ""),
                lcolor=accent, rcolor=CMAP.get(spec.get("right_accent", "TEAL"), TEAL), note=note)
    elif kind == "flow":
        flow_h(title, spec.get("steps", []), kicker=kicker, color=accent, note=note)
    elif kind == "big":
        line2 = note or " · ".join(spec.get("points", []))
        big_statement(title, line2, kicker or "KEY IDEA", color=accent)
    elif kind == "video":
        pts = []
        for item in spec.get("points", []):
            pts.append(tuple(item[:2]) if not isinstance(item, dict)
                       else (item.get("title", ""), item.get("body", "")))
        video_slide(title, spec.get("image"), spec.get("video_url"), kicker=kicker,
                    accent=accent, caption=note, points=pts or None,
                    video_embed=spec.get("video_embed"), video_file=spec.get("video_file"))
    elif kind == "image":
        image_name = spec.get("image")
        if image_name and _asset(image_name):
            if spec.get("points"):
                pts = []
                for item in spec.get("points", []):
                    pts.append(tuple(item[:2]) if not isinstance(item, dict)
                               else (item.get("title", ""), item.get("body", "")))
                img_points(title, image_name, pts, kicker=kicker, accent=accent,
                           img_w=spec.get("img_w", 7.0), note=note,
                           frame=spec.get("frame", "dark"))
            else:
                img_full(title, image_name, kicker=kicker, accent=accent, caption=note)
        else:
            content(title, spec.get("points", [note] if note else []), kicker=kicker,
                    size=spec.get("size", 17), kcolor=accent)
    elif kind == "activity":
        activity_points = spec.get("points", [])
        match = re.search(r"\bActivity\s+(\d+)\b", title, flags=re.IGNORECASE)
        activity_no = match.group(1) if match else str(spec.get("activity_no", ""))
        activity_minutes = {"1": 45, "2": 25, "3": 25, "4": 25, "5": 30,
                            "6": 30, "7": 30, "8": 15}
        activity_ka = {"1": "LO1 · A2, A4", "2": "LO2 · A5", "3": "LO2 · A5",
                       "4": "LO2 · A1, A5", "5": "LO3 · A1, A2", "6": "LO3 · A2",
                       "7": "LO3 · A1, A2", "8": "LO3 · A1, A2"}
        activity_folders = {
            "1": "activity-1-genai-agent-whatsapp",
            "2": "activity-2-excel-analysis",
            "3": "activity-3-ppt-builder",
            "4": "activity-4-tools-and-skills",
            "5": "activity-5-data-governance-policy",
            "6": "activity-6-job-redesign-role-play",
            "7": "activity-7-chatbot-security-lab",
            "8": "activity-8-security-reflection",
        }
        activity_org = {str(i): "Hands-on, no-code — ready-made website or AI agent"
                        for i in range(1, 9)}
        a = {
            "n": activity_no, "title": title,
            "minutes": spec.get("minutes", activity_minutes.get(activity_no, 45)),
            "accent": spec.get("accent", "VIOLET"),
            "org": spec.get("org", activity_org.get(activity_no, "Evidence-based classroom exercise")),
            "scenario": spec.get("scenario", note or (" · ".join(activity_points) if activity_points else
                                      "Use only synthetic data and approved lab resources.")),
            "flow": spec.get("steps", ["Scope", "Map", "Test", "Control", "Evidence"]),
            "questions": spec.get("questions", activity_points[:3] or
                                   ["Where does untrusted content meet delegated authority?",
                                    "Which hard control breaks the chain earliest?"]),
            "produce": spec.get("produce", "A documented risk decision with evidence."),
            "ka": spec.get("ka", activity_ka.get(activity_no, "K/A statements")),
            "folder": spec.get("folder", activity_folders.get(activity_no, f"activity-{activity_no}")),
        }
        activity_slide(a)
    elif kind == "break":
        brk(title or "Break", spec.get("duration", ""), color=accent)
    elif kind == "attendance":
        attendance_slide(kicker=kicker or "TRAQOM · SSG DIGITAL ATTENDANCE")
    elif kind == "lms":
        lms_slide(spec)
    elif kind == "trainer":
        is_template = "General" in title
        trainer_slide(kicker or "YOUR TRAINER",
                      spec.get("name", "Your Trainer" if is_template else C.TRAINER),
                      spec.get("role", "Complete this profile before delivery" if is_template else
                               "Principal Trainer\nTertiary Infotech Pte Ltd"),
                      spec.get("rows", [("Name", ""), ("Qualifications", ""),
                                        ("Industry experience", ""), ("Training experience", "")]
                               if is_template else
                               [("Role", "Principal Trainer"),
                                ("Expertise", "AI security, governance and enterprise IT")]),
                      spec.get("initials", "?" if is_template else "AA"), accent=accent)
    elif kind == "thankyou":
        _v30_thank_you(spec)
    else:
        content(title, spec.get("points", [note] if note else []), kicker=kicker,
                size=spec.get("size", 17), kcolor=accent)


def _add_transitions(prs, specs=None):
    """Apply restrained motion: fast fades for content, pushes for dividers."""
    from pptx.oxml import parse_xml
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    p14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for idx, slide in enumerate(prs.slides):
        spec_kind = specs[idx].get("kind", "content") if specs and idx < len(specs) else "content"
        if spec_kind == "section":
            effect = '<p:push dir="l"/>'; speed = "med"; duration = 520
        elif spec_kind == "break":
            effect = '<p:push dir="u"/>'; speed = "med"; duration = 520
        else:
            effect = '<p:fade/>'; speed = "fast"; duration = 320
        sld = slide._element
        xml = (
            f'<p:transition xmlns:p="{p_ns}" xmlns:p14="{p14}" '
            f'xmlns:a="{a_ns}" xmlns:r="{r_ns}" '
            f'p14:dur="{duration}" spd="{speed}">'
            f'{effect}'
            f'</p:transition>'
        )
        trans = parse_xml(xml)
        # Schema order for p:sld is cSld, clrMapOvr, transition, timing. Insert the
        # transition after clrMapOvr (or after cSld if clrMapOvr is absent), and
        # before any existing timing element.
        clr = sld.find(f"{{{p_ns}}}clrMapOvr")
        cSld = sld.find(f"{{{p_ns}}}cSld")
        anchor = clr if clr is not None else cSld
        idx = list(sld).index(anchor) + 1
        sld.insert(idx, trans)


def _build_v30():
    if C.VERSION.startswith("4"):
        import v40_content as V
    else:
        import v30_content as V
    global CURRENT_META
    for spec in V.SLIDES:
        if spec.get("anchor"):
            mark(spec["anchor"])
        source_labels = []
        source_details = []
        for sid in spec.get("sources", []):
            src = V.SOURCES.get(sid, {})
            title = src.get("title", "") if isinstance(src, dict) else str(src)
            url = src.get("url", "") if isinstance(src, dict) else ""
            # The footer is a compact trace key. Full titles and URLs are resolved
            # in the claim ledger; IDs alone prevent wrapping into the note/footer.
            source_labels.append(sid)
            source_details.append({"id": sid, "title": title, "url": url})
        # Real-life case studies (verified / reported) AND any slide that opts in
        # with cite=True get a VISIBLE source-citation band, so learners can see and
        # check where each claim or chart's data comes from.
        show_cite = source_details and (
            spec.get("evidence") in ("CASE-V", "CASE-R", "HIST") or spec.get("cite"))
        CURRENT_META = {"evidence": spec.get("evidence", ""), "sources": source_labels,
                        "cite": source_details if show_cite else []}
        _v30_render(spec)
    expected = getattr(V, "EXPECTED_SLIDES", 207)
    if PAGE["n"] != expected:
        raise RuntimeError(f"v3 slide count mismatch: expected {expected}, built {PAGE['n']}")
    _add_transitions(prs, V.SLIDES)
    out = os.path.join(OUTDIR, f"WSQ - Master Trainer Slides - {C.COURSE_CODE} - {C.TITLE}-v{C.VERSION.replace('.','')}.pptx")
    prs.save(out)
    with open(os.path.join(OUTDIR, "slide_map.json"), "w") as f:
        json.dump(SLIDE_MAP, f, indent=2)
    print(f"Saved: {out}")
    print(f"Slides: {PAGE['n']}")
    print("Slide map:", SLIDE_MAP)


if (C.VERSION.startswith("3") or C.VERSION.startswith("4")) and __name__ == "__main__":
    _build_v30()
    sys.exit(0)


# ================================================================ BUILD
cover()

# ---------------- ADMIN (front) ----------------
mark("admin")
section("COURSE ADMINISTRATION", "Welcome & Housekeeping", "00")
attendance_slide()
trainer_slide("YOUR TRAINER · GENERAL", "Your Trainer",
              "General Trainer template —\nto be completed by the trainer",
              [("Name", ""), ("Title / Designation", ""), ("Qualifications", ""),
               ("Areas of expertise", ""), ("Training & industry experience", ""), ("Contact", "")],
              initials="?", accent=GREY)
trainer_slide("YOUR TRAINER", C.TRAINER, "Principal Trainer\nTertiary Infotech Pte Ltd",
              [("Role", "Principal Trainer, Tertiary Infotech Pte Ltd"),
               ("Background", "PhD — 20+ years across AI, cybersecurity, data science and enterprise IT."),
               ("Delivers", "WSQ courses on generative AI, AI security, AI governance and data analytics."),
               ("Founder", "Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
              initials="AA", accent=BLUE)
content("Let's Know Each Other", [
    "Your name, organisation and role.",
    "Any generative AI or AI agent systems already live in your organisation.",
    "One AI security question you want answered by the end of Day 2.",
], kicker="ICE-BREAKER")
tile_grid("Ground Rules", [
    "Set your mobile phone to silent mode.",
    "Participate actively — no question is too small.",
    "Mutual respect: agree to disagree.",
    "One conversation at a time.",
    "Be punctual; return from breaks on time.",
    "75% attendance is required for funding.",
], kicker="HOUSEKEEPING", cols=2, size=15)
lms_slide()

table_slide("Skills Framework Alignment",
            ["TSC element", "Detail"],
            [["TSC Title", C.TSC_TITLE],
             ["TSC Code", C.TSC_CODE],
             ["Proficiency Level", C.PROFICIENCY],
             ["Knowledge statements", "K1–K5 — assessed by the Written Assessment"],
             ["Ability statements", "A1–A5 — assessed by the Case Study"]],
            kicker="WSQ ALIGNMENT", widths=[0.28, 0.72], fsize=14,
            note="This course delivers the accredited TSC through an AI security lens.")

table_slide("Learning Outcomes",
            ["LO", "Learning Outcome", "Assessed by"],
            [["LO1", C.LEARNING_OUTCOMES[0][1], "WA · Case Study"],
             ["LO2", C.LEARNING_OUTCOMES[1][1], "WA · Case Study"],
             ["LO3", C.LEARNING_OUTCOMES[2][1], "WA · Case Study"]],
            kicker="WHAT YOU WILL ACHIEVE", widths=[0.08, 0.68, 0.24], fsize=13)

table_slide("Knowledge Statements (K)",
            ["Code", "Knowledge statement"],
            [[c, d] for c, d in C.TSC_KNOWLEDGE],
            kicker="ASSESSED BY THE WRITTEN ASSESSMENT", widths=[0.09, 0.91], fsize=13,
            note="One written assessment question per knowledge statement.")

table_slide("Ability Statements (A)",
            ["Code", "Ability statement"],
            [[c, d] for c, d in C.TSC_ABILITIES],
            kicker="ASSESSED BY THE CASE STUDY", widths=[0.09, 0.91], fsize=13,
            accent=VIOLET,
            note="The case study covers every ability statement across three questions.")

two_col("Lesson Plan — Day 1", [
    ("Morning (AM attendance)", 0, True),
    ("Welcome, introductions, learning outcomes", 1),
    ("LU1 T1: The AI security threat landscape", 1),
    ("LU1 T2: Generative vs discriminative — the guardrail pattern", 1),
    ("Tea break", 1),
    ("LU1 T3: Application modes as attack surfaces", 1),
    ("Activity 1: Threat Modelling a GenAI Concierge (45 min)", 1),
    ("Lunch break 12:30–1:30pm", 1)],
    [("Afternoon (PM attendance)", 0, True),
     ("LU2 T1: Data quality, poisoning and supply chain integrity", 1),
     ("LU2 T2: Prompt injection — direct, indirect, cross-modal", 1),
     ("Tea break", 1),
     ("Activity 2: Prompt Injection & the PDPA Leak (60 min)", 1),
     ("Day 1 recap and Q&A", 1)],
    kicker="SCHEDULE · DAY 1", lhead="Morning — foundations", rhead="Afternoon — the prompt layer",
    note="8 instructional hours · 9:30am–6:30pm · 1-hour lunch · tea breaks counted within.")

two_col("Lesson Plan — Day 2", [
    ("Morning (AM attendance)", 0, True),
    ("LU2 T3: Security frameworks for GenAI and agents", 1),
    ("LU2 T4: Measuring guardrails — adversarial testing", 1),
    ("Tea break", 1),
    ("Activity 3: Selecting a Security Framework (60 min)", 1),
    ("Lunch break 12:45–1:45pm", 1)],
    [("Afternoon (PM attendance)", 0, True),
     ("LU3 T1: Agent anatomy, excessive agency, destructive execution", 1),
     ("Activity 4: Rogue Agent Post-Incident Review (60 min)", 1),
     ("LU3 T2: PDPA, PDPC and IMDA agentic governance", 1),
     ("LU3 T3: Misinformation, bias and limitations", 1),
     ("Activity 5: Agent Governance & Deployment Gate (25 min)", 1),
     ("Briefing for Assessment · WA + Case Study · TRAQOM", 1)],
    kicker="SCHEDULE · DAY 2", lhead="Morning — frameworks & measurement",
    rhead="Afternoon — agents, governance, assessment",
    note="8 instructional hours · 9:30am–6:30pm · assessment held at the end of Day 2.")

# ---------------- DAY 1 ----------------
mark("day1")
section("DAY 1 · LEARNING UNIT 1", "Foundations of AI Security", "01",
        "Why generative AI breaks the security assumptions you already rely on")

img_points("The governance gap", "diagram-incident-landscape.png",
           [("Deployed faster than governed",
             "57% of organisations run self-hosted agents; most have no inventory of them."),
            ("Prompt injection dominates",
             "The single largest source of reported GenAI and agent incidents."),
            ("Excessive agency is rising",
             "Mainstream agent deployment moved LLM03 up three places in 2026."),
            ("Monitoring is the biggest gap",
             "71% have no runtime monitoring of what their agents actually do.")],
           kicker="LU1 · T1 · WHERE WE ARE", accent=RED, img_w=6.9)

s = ncards("Why GenAI breaks classical security", D.LU1_WHY_DIFFERENT,
           kicker="LU1 · T1 · K2", accent=RED, cols=4,
           note="Classical AppSec assumes code and data are separable. A language model is the "
                "counter-example: its data IS its instructions.")

big_statement(D.BIG_STATEMENTS[0]["l1"], D.BIG_STATEMENTS[0]["l2"],
              D.BIG_STATEMENTS[0]["kicker"], color=RED)

img_full("The context window has no trust boundary", "diagram-trust-boundary.png",
         kicker="LU1 · T1 · K2", accent=RED,
         caption="Every source becomes one flat token sequence. The model cannot tell an "
                 "instruction from data — which is why LLM01 has no complete fix.")

two_col("Generative vs discriminative models",
        D.LU1_GEN_VS_DISC["left"], D.LU1_GEN_VS_DISC["right"],
        kicker="LU1 · T2 · K3", lhead="Generative — what we secure",
        rhead="Discriminative — what does the policing",
        lcolor=RED, rcolor=BLUE, note=D.LU1_GEN_VS_DISC["note"])

s = table_slide("Application modes are attack surfaces",
                D.LU1_APP_SURFACES[0], D.LU1_APP_SURFACES[1:],
                kicker="LU1 · T3 · A4", accent=VIOLET, widths=[0.17, 0.28, 0.38, 0.17],
                note="Every capability you add is a capability an attacker can aim at.")

img_full("OWASP Top 10 for LLM Applications — 2026", "diagram-owasp-llm-top10.png",
         kicker="LU1 · T3 · THE REFERENCE TAXONOMY", accent=BLUE,
         caption="Released August 2026. Ranking now weights 7,714 real incidents — "
                 "evidence, not only practitioner opinion.")

activity_slide(D.ACTIVITIES[0])
brk("Lunch Break", "12:30pm – 1:30pm")

# LU2 — prompt layer
mark("lu2")
section("DAY 1 · LEARNING UNIT 2", "Attacking and Defending the Prompt Layer", "02",
        "Injection, poisoning and the limits of prompt-based defence")

ncards("Three kinds of prompt injection",
       [(t, d.replace("\n", " ") + "\n\n→ " + n) for t, d, n, _ in D.LU2_INJECTION_TYPES],
       kicker="LU2 · T2 · K4", accent=RED, cols=3, ch_in=3.3,
       note="Indirect injection is the one that matters in production: the attacker never "
            "touches your interface, only the content your system chooses to read.")

img_points("Prompt injection in everyday work", "photo-indirect-prompt-injection-v21.png",
           D.PROMPT_INJECTION_IN_PRACTICE,
           kicker="LU2 · T2 · K4 · REALISTIC USE CASE", accent=RED, img_w=7.2,
           note="The safe design assumes the content can win and prevents it from turning into authority.")

ncards("Practice: find the hostile instruction path", D.PROMPT_PRACTICE_CASES,
       kicker="LU2 · T2 · K4 · SAFE SIMULATIONS", accent=RED, cols=4,
       note="Ask four questions: who controls the content, what can the agent reach, what action follows, and what hard control stops it?")

ncards("Why prompt-based defences fail", D.LU2_WHY_PROMPTS_FAIL,
       kicker="LU2 · T2 · K4", accent=RED, cols=4,
       note="Defensive prompting raises the attacker's cost. It never closes the hole. "
            "Treat it as friction, never as a boundary.")

table_slide("Poisoning — five vectors, one idea",
            D.LU2_POISONING[0], D.LU2_POISONING[1:],
            kicker="LU2 · T1 · K1", accent=AMBER, widths=[0.18, 0.28, 0.36, 0.18],
            note="Data quality is a security property. Whoever can write to your corpus "
                 "can write to your model's behaviour.")

activity_slide(D.ACTIVITIES[1])
content("Day 1 Recap", [
    "Prompt injection is architectural — the context window has no trust boundary.",
    "Guardrail classifiers are discriminative models policing a generative one: useful, fallible.",
    "Every application mode — summarise, infer, transform, augment — is an attack surface.",
    "Data quality is a security property: poisoning reaches training, RAG, memory and artifacts.",
    "Tomorrow: the system stops answering and starts acting.",
], kicker="DAY 1 · CLOSE", size=18)

# ---------------- DAY 2 ----------------
mark("day2")
section("DAY 2 · LEARNING UNIT 2 (cont.)", "Frameworks and Measurement", "03",
        "Choosing the right instrument, and proving the control works")

img_full("No single framework is sufficient", "diagram-framework-stack.png",
         kicker="LU2 · T3 · A3", accent=TEAL,
         caption="Each framework answers a different question. Combine them deliberately — "
                 "a threat taxonomy is not a process, and a process is not a legal obligation.")

table_slide("What each framework is — and is not",
            D.FRAMEWORK_ROLES[0], D.FRAMEWORK_ROLES[1:],
            kicker="LU2 · T3 · A3", accent=TEAL, widths=[0.30, 0.44, 0.26], fsize=12.5)

img_full("OWASP Top 10 for Agentic Applications — 2026", "diagram-owasp-asi-top10.png",
         kicker="LU2 · T3 · THE AGENT TAXONOMY", accent=VIOLET,
         caption="The agent-specific companion to the LLM list. Where the LLM list asks what the "
                 "model says, the ASI list asks what the agent does.")

formula_slide("Measuring whether a guardrail works", D.MEASUREMENT_PANELS,
              kicker="LU2 · T4 · A5", accent=BLUE,
              note="Report all three together. Any one of them alone can be gamed.")

ncards("Reading red-team results honestly", [
    ("Never read the average",
     "An aggregate attack-success rate hides the one attack family that succeeds every "
     "single time. Report per family."),
    ("Refusal rate needs a partner",
     "A falling refusal rate is not progress if the false-positive rate fell with it. "
     "You may have simply loosened the guardrail."),
    ("Vary the phrasing",
     "Test the same attack across prompt variants. A defence that holds against one "
     "phrasing is not a defence — it is a coincidence."),
    ("Re-test after every update",
     "The vendor can change model behaviour overnight. Your last red-team result "
     "expires when the model does."),
], kicker="LU2 · T4 · A5", accent=BLUE, cols=4,
   note="Publish the residual risk. A go-live gate with no stated risk tolerance is not a gate.")

activity_slide(D.ACTIVITIES[2])
brk("Lunch Break", "12:45pm – 1:45pm")

# LU3 — agents
mark("lu3")
section("DAY 2 · LEARNING UNIT 3", "Agent Autonomy, Governance and Compliance", "04",
        "When the system stops answering and starts acting")

big_statement(D.BIG_STATEMENTS[1]["l1"], D.BIG_STATEMENTS[1]["l2"],
              D.BIG_STATEMENTS[1]["kicker"], color=VIOLET)

img_full("An agent is a model plus a loop plus tools", "diagram-agent-anatomy.png",
         kicker="LU3 · T1 · K5", accent=VIOLET,
         caption="Each component adds a distinct attack surface. Autonomy is not a feature "
                 "of the model — it is a property of the architecture around it.")

table_slide("Agent capabilities and their risks",
            D.AGENT_CAPABILITY_RISK[0], D.AGENT_CAPABILITY_RISK[1:],
            kicker="LU3 · T1 · K5", accent=VIOLET, widths=[0.18, 0.28, 0.32, 0.22])

ncards("Skills and plugins are executable supply chain", D.AGENT_SUPPLY_CHAIN,
       kicker="LU3 · T1 · K5 · OWASP ASI04", accent=VIOLET, cols=4,
       note="Do not let an autonomous agent discover and install capabilities on its own. Installation is a privileged change event.")

table_slide("Trusted skill and plugin install gate",
            D.SKILL_PLUGIN_GATE[0], D.SKILL_PLUGIN_GATE[1:],
            kicker="LU3 · T1 · K5 · SUPPLY CHAIN", accent=VIOLET,
            widths=[0.20, 0.80], fsize=12.5,
            note="Approve provenance, integrity and permissions before enablement - then monitor the component at runtime.")

img_full("Agent kill chain — a real 2026 incident", "diagram-kill-chain.png",
         kicker="LU3 · T1 · K5 · CASE STUDY", accent=RED,
         caption="OpenAI → Hugging Face, July 2026. No malware, no known CVE. "
                 "The chain was assembled entirely from individually legitimate actions.")

content("What the 2026 incidents actually proved", [
    "Agents pursue goals across obstacles: a permission denial is read as a failed method, not a stop sign.",
    "Agents target other AI systems — prompt injection hidden in an HTML comment for a coding assistant.",
    "Anthropic's evaluations: a model published malicious code to PyPI that ran on 15 real systems.",
    "One model scanned ~9,000 targets and stopped only when it recognised the target was real.",
    "Replit: an agent deleted the production customer database — ASI10, in one irreversible step.",
], kicker="LU3 · T1 · K5", size=16, kcolor=RED)

img_full("Defence in depth — four layers", "diagram-defence-layers.png",
         kicker="LU3 · T1 · MICROSOFT'S MODEL", accent=TEAL,
         caption="Each layer fails differently — that is precisely why you need all four. "
                 "The application layer is the one you fully control.")

ncards("Four design patterns for safe agents", D.MS_DESIGN_PATTERNS,
       kicker="LU3 · T1 · CONTROLS", accent=TEAL, cols=4,
       note="Notice what none of these are: a better prompt. Every one is an architectural decision.")

ncards("The four-layer control stack",
       [(t, d) for t, d, _ in D.CONTROL_STACK],
       kicker="LU3 · T1 · CONTROLS", accent=BLUE, cols=4,
       note="Identity first. You cannot apply least privilege to an agent that has no distinct identity.")

activity_slide(D.ACTIVITIES[3])

# governance & compliance
big_statement(D.BIG_STATEMENTS[2]["l1"], D.BIG_STATEMENTS[2]["l2"],
              D.BIG_STATEMENTS[2]["kicker"], color=TEAL)

img_points("Human approval that actually controls the agent", "photo-human-approval-gate-v21.png",
           D.HITL_RULES, kicker="LU3 · T2 · A2 · HUMAN IN THE LOOP", accent=TEAL,
           img_w=7.1, note="Human-in-the-loop fails when the reviewer cannot see the exact action or is asked too often to think.")

img_full("Calibrating autonomy — the deterministic gate", "diagram-autonomy-gate.png",
         kicker="LU3 · T2 · A2", accent=TEAL,
         caption="IMDA: structural safeguards are preferred over prompt-based controls, and "
                 "approval checkpoints are required for irreversible actions.")

ncards("IMDA Model AI Governance for Agentic AI", D.IMDA_DIMENSIONS,
       kicker="LU3 · T2 · A2 · SINGAPORE", accent=TEAL, cols=4,
       note="Published January 2026 — the world's first governance framework written "
            "specifically for agentic AI. Some use cases are unsuitable for agents entirely.")

ncards("Six-stage organisational implementation framework", D.IMPLEMENTATION_LIFECYCLE,
       kicker="LU3 · T2 · A2 · FROM PILOT TO PRODUCTION", accent=TEAL, cols=3,
       note="NIST supplies the continuous risk cycle; IMDA sharpens the autonomy and human-accountability decisions for agents.")

ncards("Responsible AI principles become security controls", D.RESPONSIBLE_AI_SECURITY,
       kicker="LU3 · T2 · A1–A2 · RESPONSIBLE AI", accent=TEAL, cols=3,
       note="The Singapore Model AI Governance Framework and AI Verify convert principles into named owners, process checks and measurable evidence.")

table_slide("Shared responsibility across the GenAI agent lifecycle",
            D.SHARED_RESPONSIBILITY[0], D.SHARED_RESPONSIBILITY[1:],
            kicker="LU3 · T2 · A2 · ACCOUNTABILITY", accent=TEAL,
            widths=[0.25, 0.75], fsize=11.5,
            note="Responsibility is distributed, but the deploying organisation still owns the go-live decision and its real-world consequences.")

img_points("PDPA risk begins before a breach", "photo-pdpa-incident-response-v21.png",
           [("Know every data surface", "Prompts, uploads, retrieval chunks, outputs, memory, tool arguments and logs may contain personal data."),
            ("Limit use and disclosure", "Purpose, recipients, retention and overseas processing must be defined before the agent receives access."),
            ("Prepare the response", "The DPO and incident team need evidence that distinguishes human decisions, agent actions and third-party processing.")],
           kicker="LU3 · T2 · A2 · SINGAPORE", accent=RED, img_w=7.2,
           note="A model provider's contract does not remove the deploying organisation's PDPA obligations.")

img_full("PDPA accountability in the AI value chain", "diagram-pdpa-roles.png",
         kicker="LU3 · T2 · A2 · SINGAPORE", accent=RED,
         caption="PDPC GenAI guidance clarifies responsibilities across the lifecycle. The deploying organisation must evidence its own lawful, protected use.")

table_slide("PDPA obligations for an AI deployment",
            D.PDPA_DUTIES[0], D.PDPA_DUTIES[1:],
            kicker="LU3 · T2 · A2", accent=RED, widths=[0.26, 0.74], fsize=12.5)

table_slide("PDPA checklist for GenAI and agents",
            D.PDPA_AGENT_CHECKLIST[0], D.PDPA_AGENT_CHECKLIST[1:],
            kicker="LU3 · T2 · A2 · EVIDENCE", accent=RED,
            widths=[0.22, 0.78], fsize=11.5)

ncards("Bias, limits and misinformation", D.BIAS_LIMITS,
       kicker="LU3 · T3 · A1", accent=AMBER, cols=4,
       note="LLM07 rose in 2026 because a confident wrong answer that drives a tool call "
            "is no longer a quality problem — it is a security failure.")

activity_slide(D.ACTIVITIES[4])

# ---------------- CLOSE ----------------
mark("close")
section("COURSE CLOSE", "Synthesis and Assessment", "05")

ncards("Four principles to take back to work", D.CLOSING_PRINCIPLES,
       kicker="SYNTHESIS", accent=BLUE, cols=4,
       note="If you remember one thing: security for AI systems is an architecture problem "
            "wearing a prompt-engineering costume.")

table_slide("Production readiness checklist",
            D.DEPLOYMENT_CHECKLIST[0], D.DEPLOYMENT_CHECKLIST[1:],
            kicker="SYNTHESIS · ORGANISATIONAL HANDOFF", accent=BLUE,
            widths=[0.18, 0.82], fsize=11.5,
            note="No evidence means no go-live. Record exceptions as accepted residual risk with an owner and expiry date.")

content("Summary & Q&A", [
    "LU1: Generative AI breaks classical assumptions — input is executable, output drives action.",
    "LU2: Prompt injection and poisoning are architectural; defensive prompting is friction, not a boundary.",
    "LU2: Frameworks are instruments — OWASP, NIST, ATLAS, IMDA and the PDPA each answer a different question.",
    "LU3: An agent is a model plus a loop plus tools; autonomy multiplies every earlier risk.",
    "LU3: Accountability is Singapore law — the system deployer answers for what the agent does.",
], kicker="WHAT WE COVERED", size=17)

# Briefing BEFORE assessment (house rule)
content("Briefing for Assessment", [
    "Two instruments: a Written Assessment (SAQ) and a Case Study.",
    "Written Assessment: 5 short-answer questions — one for each knowledge statement K1–K5.",
    "Case Study: 3 questions — one per Learning Outcome, covering ability statements A1–A5.",
    "Format: open book. You may use the slides, the Learner Guide and your activity notes.",
    "Grading: Competent / Not Yet Competent. Re-assessment is available if you are NYC.",
    "Answer in your own words — reproducing slide text verbatim does not evidence competence.",
], kicker="BEFORE YOU BEGIN", size=17, kcolor=AMBER)

table_slide("Assessment",
            ["Instrument", "Covers", "Detail"],
            [["Written Assessment (SAQ)", "K1 – K5", C.ASSESSMENT["wa"]],
             ["Case Study", "A1 – A5 via LO1–LO3", C.ASSESSMENT["cs"]],
             ["Format", "—", C.ASSESSMENT["format"]],
             ["Grading", "—", C.ASSESSMENT["grading"]]],
            kicker="ASSESSMENT OVERVIEW", widths=[0.24, 0.20, 0.56], fsize=13)

img_full("Assessment Flow", "diagram-assessment-flow.png",
         kicker="HOW YOU WILL BE ASSESSED", accent=BLUE,
         caption="Both instruments must be completed. Your assessor confirms the outcome "
                 "with you and both parties sign off.")

content("Support", [
    "Email: enquiry@tertiaryinfotech.com",
    "Tel: +65 6318 4588",
    "Website: www.tertiarycourses.com.sg",
    "LMS/TMS: https://lms-tms.tertiaryinfotech.com",
], kicker="AFTER THE COURSE", size=18)

attendance_slide(kicker="TRAQOM · DIGITAL ATTENDANCE · END OF COURSE")


# Thank You
s = slide(); rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, SW, Inches(0.22), BLUE); rect(s, 0, Inches(7.28), SW, Inches(0.22), RED)
txt(s, 0, Inches(2.9), SW, Inches(1.3), [[("Thank You!", 54, INK, True)]], align=PP_ALIGN.CENTER)
rect(s, Inches(5.4), Inches(4.35), Inches(2.53), Inches(0.08), RED)
txt(s, 0, Inches(4.7), SW, Inches(0.8),
    [[(C.TITLE, 20, GREY, False)]], align=PP_ALIGN.CENTER)
txt(s, 0, Inches(5.25), SW, Inches(0.6),
    [[(f"{C.COURSE_CODE}  ·  Tertiary Infotech Pte Ltd  ·  UEN 201200696W", 13, GREY, False)]],
    align=PP_ALIGN.CENTER)
PAGE["n"] += 1

# ---------------- save ----------------
out = os.path.join(OUTDIR, f"WSQ - Master Trainer Slides - {C.COURSE_CODE} - {C.TITLE}-v{C.VERSION.replace('.','')}.pptx")
prs.save(out)
with open(os.path.join(OUTDIR, "slide_map.json"), "w") as f:
    json.dump(SLIDE_MAP, f, indent=2)
print(f"Saved: {out}")
print(f"Slides: {PAGE['n']}")
print("Slide map:", SLIDE_MAP)
